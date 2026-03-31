import uuid
import io
from typing import Any
import httpx
from sqlalchemy import select, text
from sqlalchemy.ext.asyncio import AsyncSession
from app.models.content_embedding import ContentEmbedding
from app.core.config import get_settings

settings = get_settings()

import os
CHUNK_SIZE = int(os.getenv("RAG_CHUNK_SIZE", "1000"))
CHUNK_OVERLAP = int(os.getenv("RAG_CHUNK_OVERLAP", "200"))


async def embed_text(text_input: str) -> list[float]:
    """Generate embedding using OpenAI text-embedding-ada-002 (1536 dims)."""
    async with httpx.AsyncClient() as client:
        response = await client.post(
            "https://api.openai.com/v1/embeddings",
            headers={"Authorization": f"Bearer {settings.OPENAI_API_KEY}"},
            json={"input": text_input, "model": "text-embedding-ada-002"},
            timeout=30.0,
        )
        response.raise_for_status()
        return response.json()["data"][0]["embedding"]


def _chunk_text(text_content: str) -> list[str]:
    if not text_content.strip():
        return []
    chunks = []
    start = 0
    while start < len(text_content):
        end = min(start + CHUNK_SIZE, len(text_content))
        chunk = text_content[start:end]
        if chunk.strip():
            chunks.append(chunk)
        if end >= len(text_content):
            break
        start = end - CHUNK_OVERLAP
        if start >= end:
            break
    return chunks


async def store_chunks(
    course_id: str, chunks: list[dict[str, Any]], db: AsyncSession
) -> None:
    for chunk in chunks:
        embedding = await embed_text(chunk["text"])
        content_embedding = ContentEmbedding(
            id=uuid.uuid4(),
            course_id=uuid.UUID(course_id),
            chunk_text=chunk["text"],
            chunk_metadata=chunk.get("metadata", {}),
            embedding=embedding,
        )
        db.add(content_embedding)
    await db.commit()


async def retrieve_relevant(
    query: str, course_id: str, db: AsyncSession, top_k: int = 5
) -> list[str]:
    query_embedding = await embed_text(query)

    result = await db.execute(
        text(
            """
            SELECT chunk_text
            FROM content_embeddings
            WHERE course_id = :course_id
            ORDER BY embedding <=> :embedding::vector
            LIMIT :limit
            """
        ),
        {
            "course_id": str(course_id),
            "embedding": str(query_embedding),
            "limit": top_k,
        },
    )
    rows = result.fetchall()
    return [row[0] for row in rows]


def extract_text_from_file(file_path: str) -> str:
    """
    Extract text from any supported file type.

    Uses the comprehensive document_extractor as primary path, which handles
    text, tables, headings, structure, etc. Falls back to basic extraction
    if the new extractor isn't available.
    """
    try:
        from app.services.document_extractor import extract_document
        doc = extract_document(file_path)
        return doc.text
    except Exception:
        pass

    # Fallback — basic extraction
    if file_path.endswith(".pdf"):
        try:
            import pypdf
            reader = pypdf.PdfReader(file_path)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except ImportError:
            raise RuntimeError("pypdf is required for PDF ingestion")
    elif file_path.endswith(".docx"):
        try:
            import docx
            doc = docx.Document(file_path)
            # Include table text alongside paragraphs
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            return "\n".join(parts)
        except ImportError:
            raise RuntimeError("python-docx is required for DOCX ingestion")
    elif file_path.endswith((".txt", ".md", ".csv")):
        with open(file_path, "r", errors="ignore") as f:
            return f.read()
    elif file_path.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                text_parts.append(" | ".join(cells))
            return "\n".join(text_parts)
        except ImportError:
            with open(file_path, "r", errors="ignore") as f:
                return f.read()
    else:
        with open(file_path, "r", errors="ignore") as f:
            return f.read()


async def ingest_document(
    course_id: str, file_path: str, db: AsyncSession
) -> int:
    """Extract content from a file and store as embedded chunks."""
    content = extract_text_from_file(file_path)

    # Use structure-aware chunking if possible
    try:
        from app.services.document_extractor import extract_document, smart_chunk
        doc = extract_document(file_path)
        text_chunks = smart_chunk(doc, max_chunk_size=1500, overlap=200)
    except Exception:
        text_chunks = _chunk_text(content)

    chunks = [
        {
            "text": chunk,
            "metadata": {
                "source": file_path,
                "chunk_index": i,
                "total_chunks": len(text_chunks),
            },
        }
        for i, chunk in enumerate(text_chunks)
    ]

    await store_chunks(course_id, chunks, db)
    return len(chunks)

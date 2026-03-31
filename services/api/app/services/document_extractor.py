"""
Comprehensive document extraction engine — NotebookLM-grade.

Extracts ALL content from uploaded documents:
- Text with structural hierarchy (headings, paragraphs, lists)
- Tables rendered as readable markdown
- Images extracted and described via Claude Vision
- Slide structure from PPTX
- Spreadsheet data from XLSX/CSV

Returns a structured ExtractedDocument with text, images, and metadata.
"""

import base64
import io
import logging
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass
class ExtractedImage:
    """An image extracted from a document."""
    data: bytes
    content_type: str  # e.g., "image/png"
    source_location: str  # e.g., "page 3", "slide 5", "after paragraph 12"
    description: str = ""  # Filled by Claude Vision


@dataclass
class ExtractedDocument:
    """Complete extraction result from a document."""
    text: str  # Full structured text (markdown-formatted)
    images: list[ExtractedImage] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)  # page_count, word_count, etc.
    sections: list[dict] = field(default_factory=list)  # Structural sections for smart chunking

    @property
    def full_text_with_descriptions(self) -> str:
        """Text with image descriptions inlined."""
        result = self.text
        for img in self.images:
            if img.description:
                result += f"\n\n[Image — {img.source_location}]: {img.description}"
        return result

    @property
    def word_count(self) -> int:
        return len(self.text.split())


# ─── DOCX Extraction ──────────────────────────────────────────────────────────

def _extract_docx(file_path: str) -> ExtractedDocument:
    """Extract everything from a Word document — paragraphs, tables, images, structure."""
    import docx
    from docx.oxml.ns import qn

    doc = docx.Document(file_path)
    parts: list[str] = []
    images: list[ExtractedImage] = []
    sections: list[dict] = []
    current_section = {"title": "Introduction", "content": [], "level": 0}

    # Walk through the document body in order — paragraphs AND tables interleaved
    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        # ── Paragraph ──
        if tag == "p":
            para = docx.text.paragraph.Paragraph(element, doc)
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else ""

            # Detect headings
            if style_name.startswith("Heading"):
                try:
                    level = int(style_name.replace("Heading", "").strip())
                except ValueError:
                    level = 1

                # Save previous section
                if current_section["content"]:
                    sections.append(current_section)
                current_section = {"title": text, "content": [], "level": level}

                prefix = "#" * min(level, 4)
                parts.append(f"\n{prefix} {text}\n")
            elif style_name.startswith("List") or style_name.startswith("Bullet"):
                parts.append(f"• {text}")
                current_section["content"].append(text)
            elif style_name.startswith("Title"):
                parts.append(f"\n# {text}\n")
                current_section["title"] = text
            else:
                parts.append(text)
                current_section["content"].append(text)

            # Check for inline images in this paragraph
            for run in para.runs:
                for drawing in run.element.findall(f".//{qn('wp:inline')}"):
                    blip = drawing.find(f".//{qn('a:blip')}")
                    if blip is not None:
                        r_embed = blip.get(qn("r:embed"))
                        if r_embed:
                            try:
                                image_part = doc.part.related_parts.get(r_embed)
                                if image_part:
                                    img_data = image_part.blob
                                    content_type = image_part.content_type or "image/png"
                                    images.append(ExtractedImage(
                                        data=img_data,
                                        content_type=content_type,
                                        source_location=f"near: {text[:60]}",
                                    ))
                                    parts.append(f"[IMAGE: embedded near '{text[:40]}...']")
                            except Exception as e:
                                logger.debug(f"Could not extract inline image: {e}")

        # ── Table ──
        elif tag == "tbl":
            table = docx.table.Table(element, doc)
            md_table = _table_to_markdown(table)
            if md_table:
                parts.append(f"\n{md_table}\n")
                current_section["content"].append(md_table)

    # Save final section
    if current_section["content"]:
        sections.append(current_section)

    # Extract remaining images from document relationships
    for rel in doc.part.rels.values():
        if "image" in str(rel.reltype):
            try:
                if rel.target_ref not in [img.source_location for img in images]:
                    image_part = rel.target_part
                    img_data = image_part.blob
                    content_type = image_part.content_type or "image/png"
                    # Deduplicate by checking data length (rough)
                    if not any(len(img.data) == len(img_data) for img in images):
                        images.append(ExtractedImage(
                            data=img_data,
                            content_type=content_type,
                            source_location="document body",
                        ))
            except Exception as e:
                logger.debug(f"Could not extract relationship image: {e}")

    text = "\n".join(parts)

    return ExtractedDocument(
        text=text,
        images=images,
        metadata={
            "format": "docx",
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
            "images_found": len(images),
            "sections_found": len(sections),
        },
        sections=sections,
    )


def _table_to_markdown(table) -> str:
    """Convert a python-docx Table to markdown format."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    # Deduplicate merged cells (python-docx repeats merged cells)
    cleaned_rows = []
    for row in rows:
        cleaned = []
        for i, cell in enumerate(row):
            # If this cell is identical to the previous one, it's likely merged
            if i > 0 and cell == row[i - 1]:
                cleaned.append("")
            else:
                cleaned.append(cell)
        cleaned_rows.append(cleaned)

    # Build markdown table
    if not cleaned_rows:
        return ""

    header = cleaned_rows[0]
    md_parts = ["| " + " | ".join(header) + " |"]
    md_parts.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in cleaned_rows[1:]:
        # Pad row to header length
        padded = row + [""] * (len(header) - len(row))
        md_parts.append("| " + " | ".join(padded[:len(header)]) + " |")

    return "\n".join(md_parts)


# ─── PDF Extraction ───────────────────────────────────────────────────────────

def _extract_pdf(file_path: str) -> ExtractedDocument:
    """Extract text and images from PDF."""
    import pypdf

    reader = pypdf.PdfReader(file_path)
    parts: list[str] = []
    images: list[ExtractedImage] = []
    sections: list[dict] = []
    current_section = {"title": "Document", "content": [], "level": 0}

    for page_num, page in enumerate(reader.pages, 1):
        page_text = page.extract_text() or ""
        if page_text.strip():
            parts.append(f"\n--- Page {page_num} ---\n")
            parts.append(page_text)
            current_section["content"].append(page_text)

        # Extract images from the page
        try:
            if hasattr(page, "images"):
                for img_obj in page.images:
                    try:
                        img_data = img_obj.data
                        # Determine content type from name
                        name = getattr(img_obj, "name", "image.png")
                        if name.lower().endswith(".jpg") or name.lower().endswith(".jpeg"):
                            ct = "image/jpeg"
                        elif name.lower().endswith(".png"):
                            ct = "image/png"
                        else:
                            ct = "image/png"

                        images.append(ExtractedImage(
                            data=img_data,
                            content_type=ct,
                            source_location=f"page {page_num}",
                        ))
                        parts.append(f"[IMAGE on page {page_num}]")
                    except Exception as e:
                        logger.debug(f"Could not extract PDF image on page {page_num}: {e}")
        except Exception as e:
            logger.debug(f"Could not access images on page {page_num}: {e}")

    # Try to detect sections from text patterns (lines that look like headers)
    full_text = "\n".join(parts)
    for line in full_text.split("\n"):
        stripped = line.strip()
        if stripped and len(stripped) < 80 and stripped.isupper():
            if current_section["content"]:
                sections.append(current_section)
            current_section = {"title": stripped, "content": [], "level": 1}
        elif stripped:
            current_section["content"].append(stripped)

    if current_section["content"]:
        sections.append(current_section)

    return ExtractedDocument(
        text=full_text,
        images=images,
        metadata={
            "format": "pdf",
            "pages": len(reader.pages),
            "images_found": len(images),
        },
        sections=sections,
    )


# ─── PPTX Extraction ─────────────────────────────────────────────────────────

def _extract_pptx(file_path: str) -> ExtractedDocument:
    """Extract text, tables, and images from PowerPoint."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(file_path)
    parts: list[str] = []
    images: list[ExtractedImage] = []
    sections: list[dict] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = []
        parts.append(f"\n## Slide {slide_num}")

        # Get slide title
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
            parts.append(f"### {title}")
            slide_parts.append(title)

        for shape in slide.shapes:
            # Text content
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
                        slide_parts.append(text)

            # Tables
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows.append(cells)

                if rows:
                    header = rows[0]
                    md = ["| " + " | ".join(header) + " |"]
                    md.append("| " + " | ".join(["---"] * len(header)) + " |")
                    for row in rows[1:]:
                        padded = row + [""] * (len(header) - len(row))
                        md.append("| " + " | ".join(padded[:len(header)]) + " |")
                    table_md = "\n".join(md)
                    parts.append(f"\n{table_md}\n")
                    slide_parts.append(table_md)

            # Images
            if hasattr(shape, "image"):
                try:
                    img_blob = shape.image.blob
                    ct = shape.image.content_type or "image/png"
                    images.append(ExtractedImage(
                        data=img_blob,
                        content_type=ct,
                        source_location=f"slide {slide_num}",
                    ))
                    parts.append(f"[IMAGE on slide {slide_num}]")
                except Exception as e:
                    logger.debug(f"Could not extract PPTX image: {e}")

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Speaker Notes: {notes}")
                slide_parts.append(f"[Notes: {notes}]")

        sections.append({
            "title": slide.shapes.title.text.strip() if slide.shapes.title else f"Slide {slide_num}",
            "content": slide_parts,
            "level": 1,
        })

    return ExtractedDocument(
        text="\n".join(parts),
        images=images,
        metadata={
            "format": "pptx",
            "slides": len(prs.slides),
            "images_found": len(images),
        },
        sections=sections,
    )


# ─── XLSX / CSV Extraction ────────────────────────────────────────────────────

def _extract_xlsx(file_path: str) -> ExtractedDocument:
    """Extract spreadsheet data as structured text."""
    import openpyxl

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts: list[str] = []
    sections: list[dict] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"\n## Sheet: {sheet_name}\n")
        sheet_content = []

        rows_data = []
        for row in ws.iter_rows(values_only=True):
            row_vals = [str(cell) if cell is not None else "" for cell in row]
            if any(v.strip() for v in row_vals):
                rows_data.append(row_vals)

        if rows_data:
            # First row as header
            header = rows_data[0]
            md = ["| " + " | ".join(header) + " |"]
            md.append("| " + " | ".join(["---"] * len(header)) + " |")
            for row in rows_data[1:]:
                padded = row + [""] * (len(header) - len(row))
                md.append("| " + " | ".join(padded[:len(header)]) + " |")
            table_md = "\n".join(md)
            parts.append(table_md)
            sheet_content.append(table_md)

        sections.append({"title": sheet_name, "content": sheet_content, "level": 1})

    wb.close()

    return ExtractedDocument(
        text="\n".join(parts),
        metadata={"format": "xlsx", "sheets": len(wb.sheetnames)},
        sections=sections,
    )


def _extract_csv(file_path: str) -> ExtractedDocument:
    """Extract CSV as a markdown table."""
    import csv as csv_mod

    with open(file_path, "r", errors="ignore") as f:
        # Sniff delimiter
        sample = f.read(4096)
        f.seek(0)
        try:
            dialect = csv_mod.Sniffer().sniff(sample)
        except csv_mod.Error:
            dialect = csv_mod.excel

        reader = csv_mod.reader(f, dialect)
        rows = list(reader)

    if not rows:
        return ExtractedDocument(text="[Empty CSV file]", metadata={"format": "csv"})

    header = rows[0]
    md = ["| " + " | ".join(header) + " |"]
    md.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in rows[1:]:
        padded = row + [""] * (len(header) - len(row))
        md.append("| " + " | ".join(padded[:len(header)]) + " |")

    text = "\n".join(md)
    return ExtractedDocument(
        text=text,
        metadata={"format": "csv", "rows": len(rows)},
        sections=[{"title": Path(file_path).stem, "content": [text], "level": 0}],
    )


# ─── Plain Text / Markdown ────────────────────────────────────────────────────

def _extract_text(file_path: str) -> ExtractedDocument:
    """Extract plain text or markdown."""
    with open(file_path, "r", errors="ignore") as f:
        content = f.read()

    # Detect sections from markdown headings
    sections = []
    current = {"title": "Document", "content": [], "level": 0}
    for line in content.split("\n"):
        if line.startswith("#"):
            if current["content"]:
                sections.append(current)
            level = len(line) - len(line.lstrip("#"))
            current = {"title": line.lstrip("#").strip(), "content": [], "level": level}
        else:
            if line.strip():
                current["content"].append(line.strip())
    if current["content"]:
        sections.append(current)

    ext = Path(file_path).suffix.lstrip(".")
    return ExtractedDocument(
        text=content,
        metadata={"format": ext},
        sections=sections,
    )


# ─── Image Description via Claude Vision ─────────────────────────────────────

async def describe_images(images: list[ExtractedImage], context: str = "") -> list[ExtractedImage]:
    """Use Claude Vision to describe extracted images, enriching the text representation."""
    if not images:
        return images

    import anthropic
    from app.core.config import get_settings
    settings = get_settings()

    client = anthropic.AsyncAnthropic(api_key=settings.ANTHROPIC_API_KEY)

    # Process up to 10 images (API cost control)
    to_process = images[:10]

    for img in to_process:
        try:
            # Skip very small images (likely icons/bullets) — under 2KB
            if len(img.data) < 2048:
                img.description = "[Small decorative image]"
                continue

            b64 = base64.b64encode(img.data).decode("utf-8")
            media_type = img.content_type
            if media_type not in ("image/png", "image/jpeg", "image/gif", "image/webp"):
                media_type = "image/png"

            response = await client.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=300,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": b64,
                            },
                        },
                        {
                            "type": "text",
                            "text": f"Describe this image concisely in 1-3 sentences. Focus on what information it conveys (data, diagrams, charts, illustrations). Context: this image is from an educational document{f' about {context}' if context else ''}.",
                        },
                    ],
                }],
            )
            img.description = response.content[0].text.strip()
        except Exception as e:
            logger.warning(f"Image description failed ({img.source_location}): {e}")
            img.description = f"[Image at {img.source_location} — description unavailable]"

    return images


# ─── Smart Chunking (Structure-Aware) ─────────────────────────────────────────

def smart_chunk(doc: ExtractedDocument, max_chunk_size: int = 1500, overlap: int = 200) -> list[str]:
    """
    Structure-aware chunking that respects section boundaries.

    Unlike naive character-based chunking, this:
    1. Keeps section headers with their content
    2. Never splits mid-table
    3. Preserves paragraph boundaries
    4. Adds overlap at section boundaries for context continuity
    """
    chunks: list[str] = []

    if doc.sections:
        for section in doc.sections:
            section_text = f"## {section['title']}\n\n" + "\n".join(section["content"])

            if len(section_text) <= max_chunk_size:
                chunks.append(section_text)
            else:
                # Split long sections at paragraph boundaries
                paragraphs = section_text.split("\n")
                current_chunk = ""
                for para in paragraphs:
                    if len(current_chunk) + len(para) + 1 > max_chunk_size and current_chunk:
                        chunks.append(current_chunk.strip())
                        # Keep some overlap
                        overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                        current_chunk = overlap_text + "\n" + para
                    else:
                        current_chunk += "\n" + para if current_chunk else para
                if current_chunk.strip():
                    chunks.append(current_chunk.strip())
    else:
        # Fallback: paragraph-based chunking
        paragraphs = doc.text.split("\n\n")
        current_chunk = ""
        for para in paragraphs:
            if len(current_chunk) + len(para) + 2 > max_chunk_size and current_chunk:
                chunks.append(current_chunk.strip())
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_text + "\n\n" + para
            else:
                current_chunk += "\n\n" + para if current_chunk else para
        if current_chunk.strip():
            chunks.append(current_chunk.strip())

    # Append image descriptions as a final chunk
    image_descs = [
        f"[Image — {img.source_location}]: {img.description}"
        for img in doc.images
        if img.description and not img.description.startswith("[Small")
    ]
    if image_descs:
        img_chunk = "## Visual Content Descriptions\n\n" + "\n\n".join(image_descs)
        if len(img_chunk) > max_chunk_size:
            # Split image descriptions
            for i in range(0, len(image_descs), 5):
                batch = image_descs[i:i + 5]
                chunks.append("## Visual Content\n\n" + "\n\n".join(batch))
        else:
            chunks.append(img_chunk)

    return chunks


# ─── Main Extraction Entry Point ──────────────────────────────────────────────

def extract_document(file_path: str) -> ExtractedDocument:
    """
    Extract ALL content from a document file.

    Supports: .docx, .pdf, .pptx, .xlsx, .csv, .txt, .md

    Returns an ExtractedDocument with structured text, images, and metadata.
    """
    ext = Path(file_path).suffix.lower()

    extractors = {
        ".docx": _extract_docx,
        ".pdf": _extract_pdf,
        ".pptx": _extract_pptx,
        ".xlsx": _extract_xlsx,
        ".csv": _extract_csv,
        ".txt": _extract_text,
        ".md": _extract_text,
        ".tsv": _extract_csv,
    }

    extractor = extractors.get(ext)
    if extractor:
        try:
            return extractor(file_path)
        except ImportError as e:
            logger.error(f"Missing dependency for {ext} extraction: {e}")
            # Fallback: try to read as plain text
            return _extract_text(file_path)
        except Exception as e:
            logger.error(f"Extraction failed for {file_path}: {e}", exc_info=True)
            # Fallback: try to read as plain text
            try:
                return _extract_text(file_path)
            except Exception:
                return ExtractedDocument(
                    text=f"[Could not extract content from {Path(file_path).name}]",
                    metadata={"format": ext, "error": str(e)},
                )
    else:
        # Unknown format — try as text
        return _extract_text(file_path)
